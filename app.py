from flask import Flask, render_template, request, send_file, url_for
from werkzeug.utils import secure_filename
import os
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'pdf', 'docx', 'xlsx', 'pptx'}

# Function to check allowed file types
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Function to extract text from images using Tesseract
def extract_text_from_image(filepath):
    image = Image.open(filepath)
    text = pytesseract.image_to_string(image)
    return text

# Function to extract text from PDF, preserving layout and formatting with fitz
def extract_text_from_pdf(filepath):
    text = ""
    with fitz.open(filepath) as pdf:
        for page_num in range(pdf.page_count):
            page = pdf[page_num]
            text += page.get_text("text")  # Keeps basic formatting
    return text

# Function to extract text from Word documents
def extract_text_from_docx(filepath):
    doc = Document(filepath)
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    return text

# Function to extract text from Excel files
def extract_text_from_excel(filepath):
    workbook = load_workbook(filepath, data_only=True)
    text = ""
    for sheet in workbook.worksheets:
        text += f"\nSheet: {sheet.title}\n"
        for row in sheet.iter_rows(values_only=True):
            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

# Function to extract text from PowerPoint files
def extract_text_from_pptx(filepath):
    presentation = Presentation(filepath)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

# Main route for handling file uploads
@app.route("/", methods=["GET", "POST"])
def upload_file():
    if request.method == "POST":
        if 'file' not in request.files:
            return "No file part"
        file = request.files['file']
        if file.filename == '':
            return "No selected file"
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Determine extraction method based on file extension
            extension = filename.rsplit('.', 1)[1].lower()
            if extension in {'jpg', 'jpeg', 'png'}:
                text = extract_text_from_image(filepath)
            elif extension == 'pdf':
                text = extract_text_from_pdf(filepath)
            elif extension == 'docx':
                text = extract_text_from_docx(filepath)
            elif extension == 'xlsx':
                text = extract_text_from_excel(filepath)
            elif extension == 'pptx':
                text = extract_text_from_pptx(filepath)
            else:
                text = "Unsupported file type."

            # Save extracted text for download if requested
            download_link = None
            if (extension == 'pdf' and len(text.splitlines()) > 50) or len(text.split()) > 500:
                download_path = os.path.join(app.config['UPLOAD_FOLDER'], 'extracted_text.txt')
                with open(download_path, 'w', encoding='utf-8') as text_file:
                    text_file.write(text)
                download_link = url_for('download_file', filename='extracted_text.txt')

            return render_template("index.html", text=text, download_link=download_link)
    return render_template("index.html")

# Route to handle file download
@app.route("/download/<filename>")
def download_file(filename):
    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
